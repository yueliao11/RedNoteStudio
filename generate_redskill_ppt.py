"""
Generate a Xiaohongshu-style pitch deck for RedNote Studio.
The deck introduces RedNote Studio as a content creation Skill for short-video platforms.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


SLIDES = [
    {
        "title": "RedNote Studio",
        "subtitle": "小红书 AI 视频笔记生成 Skill",
        "bullets": [],
        "layout": "title",
    },
    {
        "title": "创作者的痛点",
        "subtitle": "",
        "bullets": [
            "小红书需要大量短视频笔记，但写稿、配音、剪辑太耗时",
            "手里已有 PPT、讲义、产品资料，却难以快速复用",
            "多账号运营需要稳定、可批量复制的视频生产能力",
        ],
        "layout": "content",
    },
    {
        "title": "一键生成：PPT → 中文语音视频",
        "subtitle": "",
        "bullets": [
            "上传 PPTX，Agent 自动提取标题与要点",
            "自动生成自然中文讲稿，支持自定义脚本",
            "edge-tts 合成中文女声/男声神经语音",
            "ffmpeg 输出 1080P MP4，自带同步字幕",
        ],
        "layout": "content",
    },
    {
        "title": "专为小红书优化",
        "subtitle": "",
        "bullets": [
            "输出 1080P 横版视频，可裁剪为 3:4/9:16 竖版",
            "自动生成 1080×1440 小红书封面海报",
            "导出 SRT 字幕，可二次剪辑加花字",
            "结构化 metadata，方便批量管理内容资产",
        ],
        "layout": "content",
    },
    {
        "title": "三大应用场景",
        "subtitle": "",
        "bullets": [
            "知识博主：把课件/大纲批量转成口播短视频",
            "品牌种草：产品资料一键生成带货解说视频",
            "教育机构：讲义变成学生可反复观看的短视频",
        ],
        "layout": "content",
    },
    {
        "title": "为什么是 Skill，不是剪辑软件？",
        "subtitle": "",
        "bullets": [
            "剪辑软件：人工写稿、配音、剪辑，学习成本高",
            "RedNote Studio：上传 PPTX 即可，Agent 自动完成",
            "API 化设计，可嵌入任何内容工作流",
            "支持批量生产，适合矩阵号运营",
        ],
        "layout": "content",
    },
    {
        "title": "未来：内容生产 Agent",
        "subtitle": "",
        "bullets": [
            "接收选题、链接或文档，自动生成 PPT 脚本",
            "调用 RedNote Studio 生成多音色、多风格视频",
            "自动导出封面、字幕和发布文案",
            "对接小红书、抖音、B 站等内容平台",
        ],
        "layout": "content",
    },
    {
        "title": "参赛宣言",
        "subtitle": "",
        "bullets": [
            "一个真实可运行的 AI 视频笔记生成 Skill",
            "端到端：PPT 解析 → 中文讲稿 → 语音 → 视频 → 封面",
            "让每一份内容资料，都能被听见、被看见、被传播",
            "RedNote Studio：人人都是小红书视频创作者",
        ],
        "layout": "content",
    },
]


def _add_title_slide(prs, slide_data):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = title_box.text_frame
    tf.text = slide_data["title"]
    p = tf.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True

    if slide_data.get("subtitle"):
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.4), Inches(1.0))
        stf = sub_box.text_frame
        stf.text = slide_data["subtitle"]
        sp = stf.paragraphs[0]
        sp.font.size = Pt(36)


def _add_content_slide(prs, slide_data):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.6), Inches(1.2))
    tf = title_box.text_frame
    tf.text = slide_data["title"]
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True

    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(5.0))
    btf = body_box.text_frame
    btf.word_wrap = True
    for idx, bullet in enumerate(slide_data["bullets"]):
        if idx == 0:
            btf.text = bullet
            bp = btf.paragraphs[0]
        else:
            bp = btf.add_paragraph()
            bp.text = bullet
        bp.level = 0
        bp.font.size = Pt(28)
        bp.space_after = Pt(18)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        if slide_data["layout"] == "title":
            _add_title_slide(prs, slide_data)
        else:
            _add_content_slide(prs, slide_data)

    out = Path(__file__).with_name("redskill-pitchflow-deck.pptx")
    prs.save(out)
    print(f"Created REDSkill deck: {out}")


if __name__ == "__main__":
    main()
